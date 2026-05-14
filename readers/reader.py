"""Extract text from presentation/source files."""
from __future__ import annotations

import struct
import zipfile
import zlib
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree


SUPPORTED_EXTS = {
    ".pptx",
    ".pdf",
    ".docx",
    ".txt",
    ".md",
    ".xlsx",
    ".xlsm",
    ".xls",
    ".hwpx",
    ".hwp",
}


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_texts.append(line)

            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                slide_texts.append(text)
                except (AttributeError, ValueError):
                    pass

        if slide.has_notes_slide:
            note_text = slide.notes_slide.notes_text_frame.text.strip()
            if note_text:
                slide_texts.append(f"[notes] {note_text}")

        if slide_texts:
            chunks.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(chunks)


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    chunks: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if text:
            chunks.append(f"--- Page {i} ---\n{text}")
    return "\n\n".join(chunks)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        line = para.text.strip()
        if line:
            parts.append(line)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _format_cell_value(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _rows_to_text(sheet_name: str, rows: Iterable[Iterable[object]]) -> str:
    lines: list[str] = []
    for row in rows:
        cells = [_format_cell_value(cell) for cell in row]
        while cells and not cells[-1]:
            cells.pop()
        if any(cells):
            lines.append(" | ".join(cells))
    if not lines:
        return ""
    return f"--- Sheet: {sheet_name} ---\n" + "\n".join(lines)


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise RuntimeError("openpyxl is not installed; run pip install openpyxl") from e

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        chunks: list[str] = []
        for sheet in workbook.worksheets:
            text = _rows_to_text(sheet.title, sheet.iter_rows(values_only=True))
            if text:
                chunks.append(text)
        return "\n\n".join(chunks)
    finally:
        workbook.close()


def _read_xls(path: Path) -> str:
    try:
        import xlrd
    except ImportError as e:
        raise RuntimeError("xlrd is not installed; run pip install xlrd") from e

    workbook = xlrd.open_workbook(str(path))
    chunks: list[str] = []
    for sheet in workbook.sheets():
        rows = (
            (sheet.cell_value(row_index, col_index) for col_index in range(sheet.ncols))
            for row_index in range(sheet.nrows)
        )
        text = _rows_to_text(sheet.name, rows)
        if text:
            chunks.append(text)
    return "\n\n".join(chunks)


def _read_hwpx(path: Path) -> str:
    chunks: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            names = sorted(
                name
                for name in archive.namelist()
                if name.lower().startswith("contents/") and name.lower().endswith(".xml")
            )
            for name in names:
                try:
                    root = ElementTree.fromstring(archive.read(name))
                except ElementTree.ParseError:
                    continue

                texts: list[str] = []
                for elem in root.iter():
                    if elem.text and elem.text.strip():
                        texts.append(elem.text.strip())
                    if elem.tail and elem.tail.strip():
                        texts.append(elem.tail.strip())

                if texts:
                    chunks.append(f"--- Section: {Path(name).name} ---\n" + "\n".join(texts))
    except zipfile.BadZipFile as e:
        raise RuntimeError("invalid HWPX file: expected a ZIP/XML document") from e
    return "\n\n".join(chunks)


def _clean_hwp_text(text: str) -> str:
    cleaned = "".join(ch if ch in "\t\n\r" or ord(ch) >= 32 else "\n" for ch in text)
    return "\n".join(line.strip() for line in cleaned.splitlines() if line.strip())


def _extract_hwp_text_from_section(data: bytes) -> str:
    offset = 0
    parts: list[str] = []
    while offset + 4 <= len(data):
        header = struct.unpack_from("<I", data, offset)[0]
        offset += 4
        tag_id = header & 0x3FF
        size = (header >> 20) & 0xFFF
        if size == 0xFFF:
            if offset + 4 > len(data):
                break
            size = struct.unpack_from("<I", data, offset)[0]
            offset += 4
        if offset + size > len(data):
            break

        payload = data[offset : offset + size]
        offset += size
        if tag_id == 67:
            text = _clean_hwp_text(payload.decode("utf-16le", errors="ignore"))
            if text:
                parts.append(text)
    return "\n".join(parts)


def _decompress_hwp_section(data: bytes) -> bytes:
    try:
        return zlib.decompress(data, -15)
    except zlib.error as e:
        raise RuntimeError("failed to decompress HWP BodyText section") from e


def _read_hwp(path: Path) -> str:
    try:
        import olefile
    except ImportError as e:
        raise RuntimeError("olefile is not installed; run pip install olefile") from e

    if not olefile.isOleFile(str(path)):
        raise RuntimeError("invalid HWP file: expected an OLE compound document")

    chunks: list[str] = []
    with olefile.OleFileIO(str(path)) as ole:
        if not ole.exists("FileHeader"):
            raise RuntimeError("invalid HWP file: missing FileHeader stream")

        header = ole.openstream("FileHeader").read()
        if len(header) < 40:
            raise RuntimeError("invalid HWP file: incomplete FileHeader stream")
        flags = struct.unpack_from("<I", header, 36)[0]
        is_compressed = bool(flags & 0x01)
        is_encrypted = bool(flags & 0x02)
        is_distributable = bool(flags & 0x04)
        if is_encrypted:
            raise RuntimeError("encrypted HWP files are not supported")
        if is_distributable:
            raise RuntimeError("distribution-protected HWP files are not supported")

        section_paths = sorted(
            entry
            for entry in ole.listdir(streams=True, storages=False)
            if len(entry) == 2 and entry[0] == "BodyText" and entry[1].startswith("Section")
        )
        if not section_paths:
            raise RuntimeError("invalid HWP file: missing BodyText sections")

        for section_path in section_paths:
            data = ole.openstream(section_path).read()
            if is_compressed:
                data = _decompress_hwp_section(data)
            text = _extract_hwp_text_from_section(data)
            if text:
                chunks.append(f"--- Section: {section_path[-1]} ---\n{text}")
    return "\n\n".join(chunks)


def _read_text(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp949", "euc-kr"):
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def read_file(path: str | Path) -> str:
    """Read a supported file and return extracted text."""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"file not found: {p}")

    ext = p.suffix.lower()
    if ext == ".pptx":
        return _read_pptx(p)
    if ext == ".pdf":
        return _read_pdf(p)
    if ext == ".docx":
        return _read_docx(p)
    if ext in {".xlsx", ".xlsm"}:
        return _read_xlsx(p)
    if ext == ".xls":
        return _read_xls(p)
    if ext == ".hwpx":
        return _read_hwpx(p)
    if ext == ".hwp":
        return _read_hwp(p)
    if ext in {".txt", ".md"}:
        return _read_text(p)
    return ""


def read_directory(directory: str | Path) -> dict[str, str]:
    """Read every supported file in a directory."""
    d = Path(directory)
    if not d.is_dir():
        raise NotADirectoryError(f"not a directory: {d}")
    out: dict[str, str] = {}
    for p in sorted(d.iterdir()):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS:
            try:
                out[p.name] = read_file(p)
            except Exception as e:
                out[p.name] = f"[read failed: {e}]"
    return out


def iter_supported_files(directory: str | Path) -> Iterable[Path]:
    d = Path(directory)
    for p in sorted(d.iterdir()):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS:
            yield p
